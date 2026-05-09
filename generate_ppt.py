import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
# Set presentation size to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Image paths (Update these dynamically if they exist)
BG_IMG = r"C:\Users\kevin richard D\.gemini\antigravity\brain\dd7a5feb-c477-4f82-b746-c1920751240e\presentation_bg_1773937486428.png"
MOCKUP_IMG = r"C:\Users\kevin richard D\.gemini\antigravity\brain\dd7a5feb-c477-4f82-b746-c1920751240e\web_mockup_1773937885817.png"
ICON_IMG = r"C:\Users\kevin richard D\.gemini\antigravity\brain\dd7a5feb-c477-4f82-b746-c1920751240e\calendar_hero_1773938285916.png"

# Colors
NAVY_BLUE = RGBColor(18, 32, 63)
LIGHT_BLUE = RGBColor(60, 140, 210)
TEXT_GRAY = RGBColor(60, 60, 60)
WHITE = RGBColor(255, 255, 255)

slides_data = [
    {
        "title": "Appointment Booking System",
        "bullets": [
            "BCA Mini Project Presentation",
            "Submitted by: [Your Name / Roll No]",
            "Guided by: [Your Guide's Name]",
            "Academic Year: 2023-2024"
        ],
        "notes": "Good morning respected teachers. My name is [Your Name] and my roll number is [Roll No]. Today, I am going to present my BCA mini-project, which is an \"Appointment Booking System.\""
    },
    {
        "title": "Introduction",
        "bullets": [
            "A simple website to book appointments online.",
            "Replaces the old pen-and-paper booking method.",
            "Users can easily check free slots and book them.",
            "Helps doctors/admins manage their daily schedule easily.",
            "Available 24/7 on any device."
        ],
        "notes": "So, my project is a web app where people can book appointments online. Instead of calling the reception or going there in person, users can just log in and book a slot anytime they want.",
        "image": ICON_IMG
    },
    {
        "title": "Problem Statement",
        "bullets": [
            "Telephone bookings take too much time.",
            "Receptionists often make mistakes with timings.",
            "Patients have to wait in long queues at the clinic.",
            "Hard to maintain physical record books safely.",
            "No easy way to know if doctors are available today."
        ],
        "notes": "I chose this topic because I noticed the current manual system has problems. Calling takes time, and sometimes receptionists double-book times by mistake because they write it down in a hurry."
    },
    {
        "title": "Objectives of the Project",
        "bullets": [
            "To create a web-based system that is easy to use.",
            "To stop double-booking errors automatically.",
            "To let the Admin add or remove available slots.",
            "To securely save user details in a database.",
            "To save time for both users and staff."
        ],
        "notes": "My main goal with this project was to write code that stops double-booking and makes a really simple UI that anyone can use to save time.",
        "image": ICON_IMG
    },
    {
        "title": "Overview of the System",
        "bullets": [
            "Built as a dynamic website.",
            "Admin Panel: For managing the backend data.",
            "User Panel: For logging in and booking slots.",
            "Database: Stores all passwords securely and tracks bookings.",
            "Responsive: Looks good on computer and mobile screens."
        ],
        "notes": "Just to give a quick overview, the system has two main parts. A user side for the patients, and an admin side to check the daily schedules. Mobile users can also use it easily.",
        "image": MOCKUP_IMG
    },
    {
        "title": "Existing System",
        "bullets": [
            "Mostly done using phone calls or walk-ins.",
            "Staff writes down names in large registers.",
            "Changes or cancellations require another phone call.",
            "Admin has to manually check for empty time slots."
        ],
        "notes": "In the existing system, whenever we want to book a slot, we have to call. The staff looks at a notebook, finds an empty slot, and writes our name with a pen."
    },
    {
        "title": "Disadvantages of Existing System",
        "bullets": [
            "Very slow and depends totally on human effort.",
            "High chance of losing data if the register is misplaced.",
            "Annoying for users if the phone line is busy.",
            "Hard to calculate monthly reports manually."
        ],
        "notes": "The big issue here is what happens if the notebook gets lost? All patient data is gone. Also, it's very annoying to keep calling when the phone line is busy."
    },
    {
        "title": "Proposed System",
        "bullets": [
            "A complete digital shift from the old manual system.",
            "Central database keeps everything safe.",
            "Real-time slot updates (If I book 9 AM, it hides for others).",
            "Separate secure logins for Admin and Users.",
            "Search and filter options for fast access."
        ],
        "notes": "My proposed system fixes these problems. It uses a database. So, the moment one user books 10 AM, the system updates instantly, and no one else can book that same time.",
        "image": MOCKUP_IMG
    },
    {
        "title": "Advantages of Proposed System",
        "bullets": [
            "Booking can be done at midnight or from home.",
            "Reduces the receptionist's workload.",
            "100% accurate without double-booking mistakes.",
            "Admin can search any patient record in seconds.",
            "Saves paper (Eco-friendly)."
        ],
        "notes": "The biggest advantage is convenience. A user doesn't have to wait for clinic hours to call; they can book from their phone at 9 PM. It also reduces all the manual paperwork.",
        "image": ICON_IMG
    },
    {
        "title": "System Requirements",
        "bullets": [
            "Processor: Core i3 or anything basic.",
            "RAM: 4GB minimum.",
            "Storage: Just 1GB of free space is enough.",
            "Operating System: Windows 10/11.",
            "Browser: Chrome, Firefox, or MS Edge."
        ],
        "notes": "For system requirements, it's very lightweight. We don't need a heavy gaming PC. Any normal laptop with 4GB RAM and a modern web browser can run this project perfectly."
    },
    {
        "title": "Technologies Used",
        "bullets": [
            "Frontend: HTML, CSS, JavaScript (for UI design).",
            "Backend: Python with Flask framework.",
            "Database: SQLite (or MySQL).",
            "Editor Used: VS Code."
        ],
        "notes": "I used simple but powerful tools. HTML and CSS for making the pages look good, but the main brain of the app is written in Python using Flask. For storing data, I used SQLite."
    },
    {
        "title": "System Architecture",
        "bullets": [
            "Client-Server model.",
            "Browser (Client) asks for a webpage.",
            "Flask Server processes the request.",
            "Server checks the Database for information.",
            "HTML result is sent back to the user's screen."
        ],
        "notes": "The architecture is basically a client-server flow. The browser sends a normal HTTP request, my Python server runs the logic, talks to the database, and sends the page back."
    },
    {
        "title": "Modules Description",
        "bullets": [
            "User Module: Creating account, logging in, making appointments.",
            "Admin Module: Logging in securely, viewing all patient lists.",
            "Booking Module: Calendar, time slots, saving to database."
        ],
        "notes": "I broke the coding down into three main modules. The User side, the Admin dashboard, and the core Booking logic that connects the dates with the database."
    },
    {
        "title": "Step-by-Step Algorithm",
        "bullets": [
            "Step 1: Open site and Log In.",
            "Step 2: Choose the doctor and select a date.",
            "Step 3: System finds empty time slots.",
            "Step 4: Click 'Book Slot'.",
            "Step 5: System saves it in Database.",
            "Step 6: Shows 'Success' message."
        ],
        "notes": "Here is the step-by-step logic I wrote in my code. First, the user logs in, picks a date, the system queries the database for free times, and when booked, it saves it."
    },
    {
        "title": "Data Flow Diagram",
        "bullets": [
            "Explains how data moves inside the project.",
            "Shows inputs given by User (like name, date).",
            "Shows system processing (checking if slot is free).",
            "Shows data storage outputs."
        ],
        "notes": "Let me explain the DFD. It basically maps out how user input flows from the login page, through my Python server, and finally rests inside our database tables."
    },
    {
        "title": "Use Case Diagram",
        "bullets": [
            "Defines 'who does what'.",
            "The User actor can: Register, Book, View History.",
            "The Admin actor can: Login, Manage slots, Check total bookings."
        ],
        "notes": "The Use Case diagram is very simple. It shows two main actors. Users on the left booking the appointments, and the Admin on the right managing all the data."
    },
    {
        "title": "Class Diagram",
        "bullets": [
            "Represents database tables and Python objects.",
            "User Class (ID, Name, Email).",
            "Appointment Class (Date, Time, Status).",
            "Shows relationship (1 User has many Appointments)."
        ],
        "notes": "This class diagram shows the structure of my database. We have a User table, and an Appointment table. They are linked together using the user's ID as a foreign key."
    },
    {
        "title": "Sequence Diagram",
        "bullets": [
            "Shows the step-by-step timeline.",
            "User -> Browser UI -> Python Controller -> Database -> Response.",
            "Focuses on the \"Booking Appointment\" sequence."
        ],
        "notes": "The sequence diagram shows the lifeline of a single click. When a user clicks 'Submit', it goes from the HTML form to the server to the database, and back as a success message."
    },
    {
        "title": "System Testing",
        "bullets": [
            "Unit Testing: Checking small parts, like if login works.",
            "Integration Testing: Checking if Python and HTML connect properly.",
            "Validation Testing: Making sure passwords match and dates are correct.",
            "Made sure the system doesn't crash on wrong inputs."
        ],
        "notes": "Before finishing the project, I spent time testing it. I tested to make sure if a user types a wrong password or selects an old date, the system shows an error instead of breaking."
    },
    {
        "title": "Test Case Results",
        "bullets": [
            "Try booking without login -> Redirects to Login. (Pass)",
            "Enter wrong password -> Shows 'Invalid Details'. (Pass)",
            "Select a past date -> System blocks it. (Pass)",
            "Book the same time twice -> Shows 'Slot Taken'. (Pass)"
        ],
        "notes": "Here are some real test cases I tried. For example, I tried to book a slot for a day in the past. The code successfully stopped it and told me to pick a future date."
    },
    {
        "title": "Project Screenshots",
        "bullets": [
            "The Welcome / Home Page.",
            "User Registration Form.",
            "Dashboard where booking happens.",
            "Admin Panel View."
        ],
        "notes": "Let me show you what the working project looks like. Here are the screenshots. First is the clean homepage, and next is the dashboard where the user selects the date and time.",
        "image": MOCKUP_IMG
    },
    {
        "title": "Conclusion",
        "bullets": [
            "Successfully built the online booking system.",
            "It stops manual errors completely.",
            "Learned a lot about Python backend and Databases.",
            "The system is fast and easy to navigate."
        ],
        "notes": "To conclude, this project was a huge learning experience. I successfully built a working web application from scratch that actually solves the manual booking problem.",
        "image": ICON_IMG
    },
    {
        "title": "Future Enhancements",
        "bullets": [
            "Adding an online payment system (like UPI).",
            "Sending SMS or Email alerts when booked.",
            "Giving doctors their own separate login panel.",
            "Making a mobile App version."
        ],
        "notes": "Due to time limits, I couldn't add everything. But in the future, it would be great to add an online UPI payment option and send automatic SMS reminders to patients."
    },
    {
        "title": "References",
        "bullets": [
            "Official Flask Documentation (flask.palletsprojects.com)",
            "W3Schools & MDN for HTML/CSS help.",
            "StackOverflow for solving bugs.",
            "YouTube tutorials for understanding database connection."
        ],
        "notes": "Whenever I got stuck with bugs, I mostly used StackOverflow and the official Flask documentation to fix my code. I also want to thank my guide for the support."
    },
    {
        "title": "Thank You!",
        "bullets": [
            "Any Questions?"
        ],
        "notes": "That concludes my presentation. Thank you so much for your time. I am ready to answer any questions you might have."
    }
]

BLANK_SLIDE_LAYOUT = prs.slide_layouts[6] 

def add_title_slide(slide_data):
    slide = prs.slides.add_slide(BLANK_SLIDE_LAYOUT)
    
    # Add Background Image
    if os.path.exists(BG_IMG):
        slide.shapes.add_picture(BG_IMG, 0, 0, prs.slide_width, prs.slide_height)
        
        # Add a dark overlay shape
        left = top = Inches(0)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, prs.slide_width, prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        # 50% transparency is simulated by just using dark color text or python-pptx doesn't perfectly do alpha via RGBColor, 
        # so let's just make the overlay a bit smaller or just rely on the image. Actually we'll skip overlay. 
        sp = overlay._sp
        sp.getparent().remove(sp)
    
    # Add Title Box (Centered)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_SHAPE.RECTANGLE # just to avoid error
    
    p = tf.paragraphs[0]
    p.text = slide_data['title']
    p.font.bold = True
    p.font.size = Pt(64)
    p.font.name = 'Segoe UI'
    if os.path.exists(BG_IMG):
        p.font.color.rgb = WHITE
    else:
        p.font.color.rgb = NAVY_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Add Subtitle Box
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9.33), Inches(2.5))
    tf2 = txBox2.text_frame
    
    for i, bullet in enumerate(slide_data['bullets']):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(28)
        p.font.name = 'Segoe UI'
        if os.path.exists(BG_IMG):
            p.font.color.rgb = WHITE
        else:
            p.font.color.rgb = TEXT_GRAY
        p.alignment = PP_ALIGN.CENTER

    slide.notes_slide.notes_text_frame.text = slide_data['notes']

def add_content_slide(slide_data):
    slide = prs.slides.add_slide(BLANK_SLIDE_LAYOUT)
    
    # 1. Header Bar
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = NAVY_BLUE
    header_shape.line.fill.background()
    
    # Add Title text over header
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data['title']
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = 'Segoe UI'
    p.font.color.rgb = WHITE
    
    # 2. Add an accent line below header
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.2), prs.slide_width, Inches(0.05))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = LIGHT_BLUE
    accent_line.line.fill.background()
    
    # 3. Add bullets
    has_image = 'image' in slide_data and os.path.exists(slide_data['image'])
    text_width = Inches(7) if has_image else Inches(11.5)
    
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), text_width, Inches(5))
    tf = body_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(slide_data['bullets']):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = bullet
        if bullet:
            p.level = 0
        p.font.size = Pt(32)
        p.font.name = 'Segoe UI Light'
        p.font.color.rgb = TEXT_GRAY
        
        # Add space between paragraphs
        if i > 0:
            p.space_before = Pt(14)
            
    # 4. Add Image if available
    if has_image:
        # Put image on the right
        img_path = slide_data['image']
        slide.shapes.add_picture(img_path, Inches(8.2), Inches(2), height=Inches(4.5))

    # Notes
    slide.notes_slide.notes_text_frame.text = slide_data.get('notes', '')

for i, data in enumerate(slides_data):
    if i == 0 or i == len(slides_data) - 1:
        add_title_slide(data)
    else:
        add_content_slide(data)

# Add Page Numbers and Footer to all except title
for i, slide in enumerate(prs.slides):
    if i > 0 and i < len(prs.slides) - 1:
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(5), Inches(0.4))
        p = footer.text_frame.paragraphs[0]
        p.text = "Appointment Booking System Project"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        pg_num = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(1), Inches(0.4))
        p2 = pg_num.text_frame.paragraphs[0]
        p2.text = str(i + 1)
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(120, 120, 120)

prs.save('Appointment_Booking_System_Professional.pptx')
print("Professional Designer PPT generated successfully!")
